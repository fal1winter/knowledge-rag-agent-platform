"""多格式文档解析。

覆盖本地文本和 Office 格式提取，通过适配器调用 OCR 和 ASR 处理图片、音频和视频。
"""

from dataclasses import dataclass
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import Dict, Iterable, List, Protocol

from rag_agent_platform.models import DocumentAsset, ParsedDocument


class OCRClient(Protocol):
    def extract_text(self, image_uri: str) -> str:
        """返回图片 OCR 识别文本。"""


class ASRClient(Protocol):
    def transcribe(self, media_uri: str) -> str:
        """返回音视频的语音转写文本。"""


class OfficeConverter(Protocol):
    def convert(self, source_uri: str, target_extension: str) -> Path:
        """将旧版 Office 文件转换为新 XML 格式。"""


class BaseLoader(Protocol):
    supported_types: Iterable[str]

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        """解析单个资源为文本、表格、媒体转写和元数据。"""


@dataclass
class LocalOfficeConverter:
    soffice_binary: str = "soffice"

    def convert(self, source_uri: str, target_extension: str) -> Path:
        binary = shutil.which(self.soffice_binary)
        if binary is None:
            raise RuntimeError("LibreOffice soffice is required to parse legacy Office files")
        outdir = Path(tempfile.mkdtemp(prefix="office-convert-"))
        subprocess.run(
            [binary, "--headless", "--convert-to", target_extension, "--outdir", str(outdir), source_uri],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        converted = outdir / f"{Path(source_uri).stem}.{target_extension}"
        if not converted.exists():
            matches = list(outdir.glob(f"*.{target_extension}"))
            if not matches:
                raise RuntimeError(f"Office conversion did not produce .{target_extension} output")
            converted = matches[0]
        return converted


@dataclass
class PlainTextLoader:
    supported_types = {"txt", "md", "html"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        text = Path(asset.uri).read_text(encoding="utf-8", errors="ignore")
        return ParsedDocument(asset=asset, text=text, pages=[{"page": 1, "text": text}])


@dataclass
class PDFLoader:
    supported_types = {"pdf"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        import pypdf

        reader = pypdf.PdfReader(asset.uri)
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            pages.append({"page": index, "text": page_text})
        text = "\n\n".join(item["text"] for item in pages)
        return ParsedDocument(asset=asset, text=text, pages=pages)


@dataclass
class WordLoader:
    converter: OfficeConverter | None = None
    supported_types = {"doc", "docx"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        import docx

        source = Path(asset.uri)
        if asset.file_type.lower().strip(".") == "doc":
            if self.converter is None:
                raise RuntimeError("A legacy Office converter is required for .doc files")
            source = self.converter.convert(asset.uri, "docx")
        doc = docx.Document(str(source))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return ParsedDocument(asset=asset, text=text, pages=[{"page": 1, "text": text}])


@dataclass
class PowerPointLoader:
    converter: OfficeConverter | None = None
    supported_types = {"ppt", "pptx"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        from pptx import Presentation

        source = Path(asset.uri)
        if asset.file_type.lower().strip(".") == "ppt":
            if self.converter is None:
                raise RuntimeError("A legacy Office converter is required for .ppt files")
            source = self.converter.convert(asset.uri, "pptx")
        deck = Presentation(str(source))
        pages: List[Dict] = []
        for index, slide in enumerate(deck.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            pages.append({"page": index, "text": "\n".join(texts)})
        return ParsedDocument(
            asset=asset,
            text="\n\n".join(item["text"] for item in pages),
            pages=pages,
        )


@dataclass
class SpreadsheetLoader:
    supported_types = {"xls", "xlsx", "csv"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        import pandas as pd

        suffix = asset.file_type.lower().strip(".")
        frames = [pd.read_csv(asset.uri)] if suffix == "csv" else pd.read_excel(asset.uri, sheet_name=None)
        tables = []
        texts = []
        if isinstance(frames, dict):
            iterator = frames.items()
        else:
            iterator = [("sheet1", frames[0])]
        for name, frame in iterator:
            records = frame.fillna("").astype(str).to_dict(orient="records")
            tables.append({"name": name, "rows": records[:500]})
            texts.append(f"Sheet {name}\n" + frame.fillna("").astype(str).to_csv(index=False))
        return ParsedDocument(
            asset=asset,
            text="\n\n".join(texts),
            extracted_tables=tables,
        )


@dataclass
class ImageLoader:
    ocr_client: OCRClient
    supported_types = {"png", "jpg", "jpeg", "webp", "bmp", "tiff"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        text = self.ocr_client.extract_text(asset.uri)
        return ParsedDocument(asset=asset, text=text, pages=[{"page": 1, "text": text}])


@dataclass
class AudioVideoLoader:
    asr_client: ASRClient
    supported_types = {"wav", "mp3", "m4a", "mp4", "mov", "avi", "mkv"}

    def load(self, asset: DocumentAsset) -> ParsedDocument:
        transcript = self.asr_client.transcribe(asset.uri)
        return ParsedDocument(
            asset=asset,
            text=transcript,
            media_transcripts=[{"uri": asset.uri, "text": transcript}],
        )


class DocumentIngestionService:
    """按文件类型分发到对应的解析器。"""

    def __init__(self, loaders: List[BaseLoader]):
        self.loaders = loaders

    def parse(self, asset: DocumentAsset) -> ParsedDocument:
        file_type = asset.file_type.lower().strip(".")
        for loader in self.loaders:
            if file_type in set(loader.supported_types):
                return loader.load(asset)
        raise ValueError(f"Unsupported document type: {asset.file_type}")
